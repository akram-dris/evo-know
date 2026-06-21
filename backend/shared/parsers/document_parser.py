import os
from pathlib import Path
from PyPDF2 import PdfReader
import docx

class DocumentParser:
    """Parses PDF, DOCX, DOC, PPTX, and TXT files into raw text."""

    def parse(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return self._parse_pdf(file_path)
        elif ext == '.docx':
            return self._parse_docx(file_path)
        elif ext == '.doc':
            return self._parse_doc(file_path)
        elif ext == '.pptx':
            return self._parse_pptx(file_path)
        elif ext == '.txt':
            return self._parse_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _parse_pdf(self, file_path: str) -> str:
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text.append(t)
        return "\n".join(text)

    def _parse_docx(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text])

    def _parse_doc(self, file_path: str) -> str:
        import subprocess
        try:
            result = subprocess.run(["antiword", file_path], capture_output=True, text=True, encoding='utf-8', errors='ignore')
            if result.returncode == 0:
                return result.stdout
        except Exception as e:
            pass
        raise ValueError("Le format .doc nécessite l'outil system 'antiword' sur le serveur. Veuillez convertir le fichier en .docx.")

    def _parse_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            raise ValueError(f"Failed to parse .pptx file: {e}")

    def _parse_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
